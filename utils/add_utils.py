"""
Import necessary libraries
"""

import os
import base64
import sys
sys.path.append("./")
from utils.log_utils import logger
import streamlit as st
import PyPDF2
import pdfplumber
import docx
from pptx import Presentation
from PIL import Image
import pytesseract

def get_base64_image(image_path):
    with open(image_path, "rb") as img_file:
        image_base64 = base64.b64encode(img_file.read()).decode("utf-8")
        if image_base64:
            logger.info(f"Rendered image : {image_path}")
    return image_base64

def set_bg(main_bg):
    '''
    A function to unpack an image from root folder and set as bg.

    Returns
    -------
    The background.
    '''
    # set bg name
    main_bg_ext = "png"
        
    st.markdown(
        f"""
        <style>
        .stApp {{
            background: url(data:image/{main_bg_ext};base64,{base64.b64encode(open(main_bg, "rb").read()).decode()});
            background-size: cover
        }}
        </style>
        """,
        unsafe_allow_html=True
    )


def render_chat_message(role, avatar_path, content):
    with open(avatar_path, "rb") as img_file:
        avatar_base64 = base64.b64encode(img_file.read()).decode("utf-8")
        alignment = "margin-left: auto;" if role == "user" else "margin-right: auto;"
        bubble_bg = "#004080" if role == "user" else "#FFFFFF"
        bubble_text = "#FFFFFF" if role == "user" else "black"
        st.markdown(f"""
    <div style="{alignment} max-width: 75%; margin-bottom: 4px;">
    <div style="display: flex; align-items: center; gap: 8px; flex-direction: {'row-reverse' if role == 'user' else 'row'};">
    <img src="data:image/png;base64,{avatar_base64}"
                        style="width: 47px; height: 47px; border-radius: 60%; flex-shrink: 0;">
    <div style="background: {bubble_bg}; color:{bubble_text};padding: 12px; border-radius: 16px; line-height: 1.4;">
                    {content}
    </div>
    </div>
    </div>
    """, unsafe_allow_html=True)
        
# Helper functions
def read_txt(uploaded):
    return uploaded.read().decode('utf-8', errors='ignore')

def read_pdf(uploaded):
    text = []
    with pdfplumber.open(uploaded) as pdf:
        for page in pdf.pages:
            text.append(page.extract_text() or "")
    return "\n\n".join(text)

def read_docx(uploaded):
    doc = docx.Document(uploaded)
    return "\n\n".join(para.text for para in doc.paragraphs)

def read_pptx(uploaded):
    prs = Presentation(uploaded)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n\n".join(text_runs)

def read_image(uploaded):
    img = Image.open(uploaded)
    return pytesseract.image_to_string(img)